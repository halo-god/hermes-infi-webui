"""File-safety helpers: bounded upload reads + path-traversal confinement."""
from __future__ import annotations

import asyncio
import base64
import logging
import os
import uuid
from dataclasses import dataclass
from functools import partial

from fastapi import HTTPException, UploadFile

logger = logging.getLogger(__name__)

_UPLOAD_CHUNK = 1024 * 1024  # 1 MiB


def extract_pdf_text(data: bytes) -> str | None:
    """Extract plain text from PDF bytes using PyMuPDF (fitz)."""
    try:
        import pymupdf
        doc = pymupdf.open(stream=data, filetype="pdf")
        parts: list[str] = []
        for page in doc:
            txt = page.get_text()
            if txt:
                parts.append(txt)
        return "\n\n".join(parts) if parts else None
    except Exception:
        return None


_DOCX_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
    "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
}


def _extract_images_from_drawing(
    drawing_elem,
    image_map: dict[str, str],
    text_runs: list[str],
) -> None:
    """Find all <a:blip r:embed="rId..."> inside a <w:drawing> and append
    base64 <img> tags to text_runs."""
    for blip in drawing_elem.iter():
        tag = blip.tag.split("}")[-1] if "}" in blip.tag else blip.tag
        if tag != "blip":
            continue
        embed_attr = None
        for attr_name, attr_val in blip.attrib.items():
            if attr_name.endswith("}embed") or attr_name == "embed":
                embed_attr = attr_val
                break
        if embed_attr and embed_attr in image_map:
            text_runs.append(
                f'<img src="{image_map[embed_attr]}" style="max-width:100%;height:auto;border-radius:6px;margin:6px 0" />'
            )


def extract_docx_html(raw: bytes) -> str | None:
    """Convert a .docx file to sanitized preview HTML (mammoth engine).

    Mammoth produces semantic HTML (headings, lists, tables, footnotes) from
    the OOXML — far better structure than the hand-rolled parser for complex
    documents. Images are inlined as base64 data URIs (sized-capped, like the
    pptx extractor). Mammoth escapes all text content itself; the surrounding
    try/except keeps malformed files falling through to None.
    """
    try:
        import base64
        import io

        import mammoth

        _MAMMOTH_IMAGE_MAX = 2 * 1024 * 1024

        def _convert_image(image):
            with image.open() as image_bytes:
                data = image_bytes.read()
            if len(data) > _MAMMOTH_IMAGE_MAX:
                return {}  # drop oversized images from the preview
            ctype = image.content_type or "image/png"
            return {"src": f"data:{ctype};base64," + base64.b64encode(data).decode("ascii")}

        result = mammoth.convert_to_html(
            io.BytesIO(raw),
            convert_image=_convert_image,
        )
        html = (result.value or "").strip()
        if not html:
            return "<p><em>(空文档)</em></p>"
        return html
    except Exception:
        return None


_XLSX_MAX_ROWS = 500
_XLSX_MAX_COLS = 50
# Workbooks below this size are loaded in full-fidelity mode (column widths,
# merges, fills, fonts, number formats, freeze panes). Larger ones fall back
# to streaming read_only mode: plain grid + header row only.
_XLSX_RICH_MAX_BYTES = 5 * 1024 * 1024

# Spreadsheet preview styling. The extractor output is rendered inside a
# sandboxed srcdoc iframe which CANNOT see app CSS — so the styles must be
# inlined in a complete HTML document for xlsx/csv previews to look like a
# spreadsheet instead of a bare unstyled table.
_SPREADSHEET_CSS = """
body { margin: 0; padding: 14px 16px; background: #fff; color: #1f2328;
       font: 13px/1.5 -apple-system, "PingFang SC", "Microsoft YaHei", sans-serif; }
h3.sheet { font-size: 14px; margin: 14px 0 8px; padding-bottom: 4px;
           border-bottom: 2px solid #b8852a; color: #444; }
h3.sheet:first-child { margin-top: 0; }
table { border-collapse: collapse; margin: 0 0 12px; }
th, td { border: 1px solid #e2e4e8; padding: 5px 10px; text-align: left;
         white-space: nowrap; min-width: 48px; }
th { background: #f5f6f7; font-weight: 600; position: sticky; top: 0;
     box-shadow: inset 0 -1px 0 #d8dade; }
tbody tr:nth-child(even) { background: #fafbfc; }
td.num { text-align: right; font-variant-numeric: tabular-nums; }
td.fc, th.fc { position: sticky; left: 0; background: #fff; z-index: 1;
               box-shadow: inset -1px 0 0 #e2e4e8; }
th.fc { z-index: 2; }
tbody tr:nth-child(even) td.fc { background: #fafbfc; }
p.note { color: #8a8f98; font-size: 12px; margin: 4px 0 12px; }
"""


def _spreadsheet_doc(body_html: str) -> str:
    """Wrap extracted table markup in a complete styled HTML document."""
    return (
        "<!DOCTYPE html><html><head><meta charset=\"utf-8\">"
        f"<style>{_SPREADSHEET_CSS}</style></head><body>{body_html}</body></html>"
    )


def is_legacy_spreadsheet_html(kind: str, content: str | None) -> bool:
    """True when a stored xlsx/csv preview predates the styled-document
    format (bare table fragment, no DOCTYPE) and should be re-extracted."""
    return kind in ("xlsx", "csv") and bool(content) and "<!DOCTYPE" not in content[:200]


def _xlsx_fmt_value(value, fmt: str | None) -> str:
    """Render a cell value the way Excel's number format intends."""
    import datetime
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, datetime.datetime):
        return value.strftime("%Y-%m-%d %H:%M")
    if isinstance(value, datetime.date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, (int, float)):
        f = fmt or ""
        if "%" in f:
            decimals = len(f.split(".")[1].rstrip("%")) if "." in f else 0
            return f"{value * 100:.{decimals}f}%"
        if "#,##0" in f:
            decimals = len(f.split(".")[1]) if "." in f else 0
            return f"{value:,.{decimals}f}"
        if f in ("0.00", "0.000"):
            return f"{value:.{len(f) - 2}f}"
        # General / integers: drop float artifacts from data_only values.
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value)
    return str(value)


def _xlsx_inline_style(cell) -> str:
    """Translate openpyxl cell styling to inline CSS (best effort)."""
    styles: list[str] = []
    try:
        fill = cell.fill
        if fill is not None and fill.patternType == "solid":
            rgb = getattr(fill.fgColor, "rgb", None)
            if isinstance(rgb, str) and len(rgb) >= 6:
                styles.append(f"background:#{rgb[-6:]}")
    except Exception:
        pass
    try:
        font = cell.font
        if font is not None:
            if font.bold:
                styles.append("font-weight:700")
            if font.italic:
                styles.append("font-style:italic")
            color = getattr(font.color, "rgb", None) if font.color else None
            if isinstance(color, str) and len(color) >= 6 and color[-6:] != "000000":
                styles.append(f"color:#{color[-6:]}")
    except Exception:
        pass
    try:
        align = cell.alignment
        if align is not None:
            if align.horizontal == "center":
                styles.append("text-align:center")
            elif align.horizontal == "right":
                styles.append("text-align:right")
            if align.wrap_text:
                styles.append("white-space:normal")
    except Exception:
        pass
    return ";".join(styles)


def _xlsx_rich_sheet_html(ws) -> str:
    """Full-fidelity sheet render: widths, heights, merges, fills, fonts,
    number formats, freeze panes, alignment."""
    from html import escape
    from openpyxl.utils import get_column_letter

    # Merged ranges: anchor -> (rowspan, colspan); covered coords get skipped.
    anchors: dict[tuple[int, int], tuple[int, int]] = {}
    covered: set[tuple[int, int]] = set()
    try:
        for rng in ws.merged_cells.ranges:
            anchors[(rng.min_row, rng.min_col)] = (rng.max_row - rng.min_row + 1,
                                                   rng.max_col - rng.min_col + 1)
            for r in range(rng.min_row, rng.max_row + 1):
                for c in range(rng.min_col, rng.max_col + 1):
                    if (r, c) != (rng.min_row, rng.min_col):
                        covered.add((r, c))
    except Exception:
        pass

    # Column widths (Excel width unit ≈ 7px per char + padding).
    col_tags: list[str] = []
    try:
        for idx in range(1, min(ws.max_column, _XLSX_MAX_COLS) + 1):
            dim = ws.column_dimensions.get(get_column_letter(idx))
            w = dim.width if dim is not None and dim.width else None
            if w:
                col_tags.append(f'<col style="width:{round(w * 7 + 5)}px">')
            else:
                col_tags.append("<col>")
    except Exception:
        col_tags = []
    colgroup = f"<colgroup>{''.join(col_tags)}</colgroup>" if col_tags else ""

    # Freeze panes: sticky first column when frozen at column ≥ 2. (A frozen
    # first row is equivalent to the sticky <thead> already in the base CSS.)
    frozen_col = False
    try:
        if ws.freeze_panes and ws.freeze_panes not in ("A1", None):
            from openpyxl.utils.cell import coordinate_from_string
            _c, _r = coordinate_from_string(ws.freeze_panes)
            from openpyxl.utils import column_index_from_string
            frozen_col = column_index_from_string(_c) >= 2
    except Exception:
        pass

    thead_rows: list[str] = []
    body_rows: list[str] = []
    row_count = 0
    truncated_cols = False
    for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, _XLSX_MAX_ROWS),
                            max_col=min(ws.max_column, _XLSX_MAX_COLS)):
        if ws.max_column > _XLSX_MAX_COLS:
            truncated_cols = True
        row_idx = row[0].row if row else row_count + 1
        cells_html: list[str] = []
        height = None
        try:
            rd = ws.row_dimensions.get(row_idx)
            if rd is not None and rd.height:
                height = round(rd.height)
        except Exception:
            pass
        tr_style = f' style="height:{height}px"' if height else ""
        for cell in row:
            pos = (cell.row, cell.column)
            if pos in covered:
                continue
            span = anchors.get(pos)
            span_attr = ""
            if span:
                rs, cs = span
                span_attr = (f' rowspan="{rs}" colspan="{cs}"' if rs > 1 or cs > 1 else "")
            tag = "th" if row_idx == 1 else "td"
            classes: list[str] = []
            if frozen_col and cell.column == 1:
                classes.append("fc")
            if (isinstance(cell.value, (int, float)) and not isinstance(cell.value, bool)
                    and not (cell.alignment and cell.alignment.horizontal)):
                classes.append("num")
            cls = f' class="{" ".join(classes)}"' if classes else ""
            style = _xlsx_inline_style(cell)
            style_attr = f' style="{style}"' if style else ""
            text = _xlsx_fmt_value(cell.value, cell.number_format) if cell.value is not None else ""
            cells_html.append(f"<{tag}{cls}{style_attr}{span_attr}>{escape(text)}</{tag}>")
        tr = f'<tr{tr_style}>{"".join(cells_html)}</tr>'
        if row_idx == 1:
            thead_rows.append(tr)
        else:
            body_rows.append(tr)
        row_count += 1

    if ws.max_row > _XLSX_MAX_ROWS:
        truncated_note = f'<p class="note">(仅显示前 {_XLSX_MAX_ROWS} 行，已截断)</p>'
    else:
        truncated_note = ""
    col_note = f'<p class="note">(仅显示前 {_XLSX_MAX_COLS} 列，已截断)</p>' if truncated_cols else ""
    table = (
        "<table>" + colgroup
        + (f'<thead>{"".join(thead_rows)}</thead>' if thead_rows else "")
        + f'<tbody>{"".join(body_rows)}</tbody>'
        + "</table>" + truncated_note + col_note
    )
    return table


def extract_xlsx_html(raw: bytes) -> str | None:
    """Convert an .xlsx workbook to a styled spreadsheet-like preview.

    Workbooks under _XLSX_RICH_MAX_BYTES are rendered in full fidelity
    (column widths, row heights, merged cells, fills, font colors, number
    formats, freeze panes, alignment); larger ones stream in read_only mode
    with the base grid styling.
    """
    try:
        import io
        from html import escape

        from openpyxl import load_workbook

        # Validate: .xlsx must be a ZIP archive (PK\x03\x04 magic).
        if len(raw) < 4 or raw[:4] != b"PK\x03\x04":
            # If it looks like HTML/text, the agent likely wrote preview markup
            # into a .xlsx filename instead of the real binary. Surface this
            # clearly so the user / agent knows what went wrong.
            preview = raw[:200].decode("utf-8", "ignore").strip()
            return _spreadsheet_doc(
                '<p><em style="color:#c0392b">⚠ 文件格式错误：该文件不是有效的 .xlsx 工作簿。</em></p>'
                '<p>可能原因：AI 助手在生成文件时将 HTML 预览写入了 <code>.xlsx</code> 扩展名，'
                '而非真正的 Excel 二进制内容。</p>'
                '<p>解决方法：请使用 <code>write_file</code> 工具生成 <code>.md</code> 或 <code>.txt</code> 文件，'
                '或者上传真实的 .xlsx 文件。</p>'
                f'<pre style="background:#f8f9fa;padding:8px;border-radius:4px;font-size:12px;white-space:pre-wrap">{escape(preview)}</pre>'
            )

        rich = len(raw) <= _XLSX_RICH_MAX_BYTES
        wb = load_workbook(io.BytesIO(raw), read_only=not rich, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f'<h3 class="sheet">{escape(ws.title)}</h3>')
            if rich:
                parts.append(_xlsx_rich_sheet_html(ws))
            else:
                parts.append(_xlsx_stream_sheet_html(ws))
        wb.close()
        return _spreadsheet_doc("\n".join(parts)) if parts else _spreadsheet_doc("<p><em>(空工作簿)</em></p>")
    except Exception:
        return None


def _xlsx_stream_sheet_html(ws) -> str:
    """Streaming (read_only) sheet render: plain grid + header row."""
    from html import escape

    rows_html: list[str] = []
    header_done = False
    truncated_cols = False
    row_count = 0
    for row in ws.iter_rows():
        if row_count >= _XLSX_MAX_ROWS:
            rows_html.append(f'</tbody><p class="note">(仅显示前 {_XLSX_MAX_ROWS} 行，已截断)</p>')
            break
        cells = row[:_XLSX_MAX_COLS]
        if len(row) > _XLSX_MAX_COLS:
            truncated_cols = True
        if not header_done:
            cell_html = "".join(
                f"<th>{escape(str(c.value)) if c.value is not None else ''}</th>" for c in cells
            )
            rows_html.append(f"<thead><tr>{cell_html}</tr></thead><tbody>")
            header_done = True
        else:
            cell_html = "".join(
                (
                    '<td class="num">' if isinstance(c.value, (int, float)) else "<td>"
                ) + (escape(str(c.value)) if c.value is not None else "") + "</td>"
                for c in cells
            )
            rows_html.append(f"<tr>{cell_html}</tr>")
        row_count += 1
    table = "<table>" + "".join(rows_html) + ("</tbody>" if header_done else "") + "</table>"
    if truncated_cols:
        table += f'<p class="note">(仅显示前 {_XLSX_MAX_COLS} 列，已截断)</p>'
    return table


def extract_pptx_html(raw: bytes) -> str | None:
    """Convert a .pptx presentation's slide text to sanitized preview HTML.

    Enhanced vs the original text-only extractor:
      - title/centered placeholders render as headings (h3/h4), other text as <p>
      - tables (GraphicFrame) render as real HTML <table>
      - embedded images render inline as base64 <img> (fallback to a
        "图片" caption marker when encoding fails)
      - slide ordering preserved, sanitized via html.escape
    """
    try:
        import io
        from html import escape

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(raw))
        parts: list[str] = []
        inline_img_count = 0
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f'<div class="slide"><h4>Slide {i}</h4>')
            for shape in slide.shapes:
                # ── Tables ──
                if shape.has_table:
                    tbl = shape.table
                    parts.append("<table>")
                    for ri, row in enumerate(tbl.rows):
                        tag = "th" if ri == 0 else "td"
                        cells = "".join(
                            f"<{tag}>{escape(cell.text or '')}</{tag}>"
                            for cell in row.cells
                        )
                        parts.append(f"<tr>{cells}</tr>")
                    parts.append("</table>")
                    continue

                # ── Images / pictures ──
                if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP):
                    # Cap inline images — a picture-heavy deck would otherwise
                    # balloon the preview HTML into multi-MB base64 and stall
                    # the browser.
                    if inline_img_count < _PPTX_MAX_INLINE_IMAGES:
                        img_b64, img_mime = _pptx_shape_image_b64(shape)
                        if img_b64:
                            inline_img_count += 1
                            parts.append(
                                f'<div class="slide-img"><img src="data:{img_mime};base64,'
                                f'{img_b64}" alt="图片" /></div>'
                            )
                    continue

                # ── Text frames (title / body / notes) ──
                if not shape.has_text_frame:
                    continue
                # Title placeholder (or any shape flagged as title)
                is_title = (
                    getattr(shape, "is_placeholder", False)
                    and getattr(shape, "placeholder_format", None) is not None
                    and getattr(shape.placeholder_format, "type", None) is not None
                    and "TITLE" in str(getattr(shape.placeholder_format, "type", "")).upper()
                )
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text or "" for run in para.runs) or (para.text or "")
                    if not text.strip():
                        continue
                    if is_title:
                        parts.append(f"<h3>{escape(text.strip())}</h3>")
                    else:
                        parts.append(f"<p>{escape(text.strip())}</p>")
            parts.append("</div>")
        return "\n".join(parts) if parts else "<p><em>(空演示文稿)</em></p>"
    except Exception:
        return None


# Inline images bigger than this are dropped from the preview (base64 expands
# by ~33% — a multi-MB picture would balloon the preview HTML for no benefit).
_PPTX_IMAGE_MAX_BYTES = 2 * 1024 * 1024
# Total inline images per deck — picture-heavy decks (hundreds of images)
# would otherwise produce multi-MB preview HTML and stall the browser.
_PPTX_MAX_INLINE_IMAGES = 40


def _pptx_shape_image_b64(shape) -> tuple[str, str] | tuple[None, str]:
    """Extract an inline base64 + content-type from a pptx picture/group shape.

    Returns (b64, mime); mime is the picture's real content type (image/jpeg
    for JPEGs — hardcoding image/png breaks those). Oversized images yield
    (None, "") so the preview stays light. For grouped shapes only the first
    picture is taken (preview compromise).
    """
    try:
        import base64 as _b64

        from pptx.enum.shapes import MSO_SHAPE_TYPE

        def _walk(s):
            if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE and hasattr(s, "image"):
                try:
                    blob = s.image.blob
                    if len(blob) > _PPTX_IMAGE_MAX_BYTES:
                        return None
                    return (blob, s.image.content_type or "image/png")
                except Exception:
                    return None
            if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                for sub in getattr(s, "shapes", []):
                    hit = _walk(sub)
                    if hit:
                        return hit
            return None

        hit = _walk(shape)
        if hit:
            blob, mime = hit
            return (_b64.b64encode(blob).decode("ascii"), mime)
    except Exception:
        return None, ""
    return None, ""


def extract_csv_html(raw: bytes) -> str | None:
    """Convert CSV bytes to a styled HTML table preview (first 200 rows)."""
    try:
        import io
        import csv
        from html import escape

        text = raw.decode("utf-8", "ignore")
        reader = csv.reader(io.StringIO(text))
        parts: list[str] = ["<table>"]
        row_count = 0
        for row in reader:
            if row_count >= 200:
                parts.append('<p class="note">(仅显示前 200 行，已截断)</p>')
                break
            tag = "th" if row_count == 0 else "td"
            cells = "".join(f"<{tag}>{escape(c)}</{tag}>" for c in row)
            parts.append(f"<tr>{cells}</tr>")
            row_count += 1
        parts.append("</table>")
        return _spreadsheet_doc("\n".join(parts)) if row_count > 0 else None
    except Exception:
        return None


def extract_rtf_html(raw: bytes) -> str | None:
    """Extract plain text from RTF and wrap in <p> tags."""
    try:
        import re
        text = raw.decode("utf-8", "ignore")
        # Strip RTF control words and groups
        text = re.sub(r"\\'[0-9a-fA-F]{2}", "", text)
        text = re.sub(r"\\[a-zA-Z]+-?\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        text = re.sub(r"\\\*", "", text)
        text = re.sub(r"\\\n", "\n", text)
        text = text.strip()
        if not text:
            return None
        from html import escape
        paragraphs = [f"<p>{escape(p)}</p>" for p in text.split("\n\n") if p.strip()]
        return "\n".join(paragraphs) if paragraphs else None
    except Exception:
        return None


OFFICE_EXTRACTORS = {
    "docx": extract_docx_html,
    "xlsx": extract_xlsx_html,
    "pptx": extract_pptx_html,
    "csv": extract_csv_html,
    "rtf": extract_rtf_html,
    # Legacy OLE2 formats (.doc/.xls/.ppt) and OpenDocument (.odt/.ods/.odp)
    # — require LibreOffice (soffice) at runtime; absent soffice yields a
    # clear "convert to docx/xlsx" hint instead of an unreadable file.
    # Registered below via partial (extractor signature stays single-arg).
    "doc": None,
    "xls": None,
    "ppt": None,
    "odt": None,
    "ods": None,
    "odp": None,
}

from app.core.office_legacy import extract_legacy_doc_html  # noqa: E402

for _ext in ("doc", "xls", "ppt", "odt", "ods", "odp"):
    OFFICE_EXTRACTORS[_ext] = partial(extract_legacy_doc_html, ext=_ext)

# Formats that get the unified LibreOffice PDF preview. csv/rtf are excluded:
# their structured extraction (HTML table / text) is strictly better for both
# preview and AI injection, and PDF conversion would waste a soffice run.
_PDF_PREVIEW_EXTS = frozenset({
    "docx", "xlsx", "pptx", "doc", "xls", "ppt", "odt", "ods", "odp",
})



PLAIN_TEXT_EXTS = frozenset({
    "md", "txt", "json", "csv", "html", "htm", "js", "ts", "py", "go", "rs",
    "yaml", "yml", "toml", "sh", "bash", "log", "xml", "css", "diff", "patch",
})


def is_text_extractable(kind: str) -> bool:
    """Return True for file kinds we can extract human-readable text from."""
    kind = kind.lower()
    return kind in PLAIN_TEXT_EXTS or kind == "pdf" or kind in OFFICE_EXTRACTORS


@dataclass
class ProcessedUpload:
    content: str | None
    storage_key: str | None
    size_bytes: int
    # Object-storage key of the LibreOffice-converted PDF (Office uploads
    # only); None when soffice is missing/failed or the format isn't Office.
    preview_pdf_key: str | None = None


def _decode_text(raw: bytes) -> str:
    """Decode bytes to text using charset detection (charset-normalizer) so
    GBK/Big5/Shift-JIS files aren't silently truncated by utf-8 ignore.

    Short Chinese text often misfires as Korean cp949 (overlapping byte
    ranges); we retry constrained to Chinese encodings when that happens."""
    try:
        from charset_normalizer import from_bytes
        result = from_bytes(raw).best()
        if result is not None:
            enc = (result.encoding or "").lower()
            if enc in ("cp949", "euc-kr", "iso-2022-kr"):
                cn = from_bytes(raw, cp_isolation=["utf-8", "gb18030", "gbk", "big5"]).best()
                if cn is not None:
                    return str(cn)
            return str(result)
    except Exception:  # noqa: BLE001 — fall back to utf-8 ignore
        pass
    return raw.decode("utf-8", "ignore")


def _strip_exif(raw: bytes, ext: str) -> bytes:
    """P2-file: strip EXIF metadata from images (GPS, camera, timestamps).
    Returns the original bytes if Pillow is missing or the file isn't an
    image / can't be processed — never blocks the upload."""
    if ext.lower().lstrip(".") not in ("jpg", "jpeg", "png", "webp"):
        return raw
    try:
        from io import BytesIO
        from PIL import Image
        img = Image.open(BytesIO(raw))
        # Re-create the image without metadata. Using the pixel data copy
        # approach (not getdata which is deprecated in Pillow 14).
        cleaned = Image.new(img.mode, img.size)
        cleaned.paste(img)
        buf = BytesIO()
        # Preserve format; PNG has no EXIF so this is mainly for JPEG.
        fmt = "PNG" if ext.lower().lstrip(".") == "png" else "JPEG"
        cleaned.save(buf, format=fmt)
        return buf.getvalue()
    except Exception:  # noqa: BLE001
        return raw


def _extract_archive(raw: bytes, ext: str) -> str | None:
    """P2-file: extract text from a zip/tar/gz archive by recursively
    processing each member. Returns a single concatenated text with file
    separators, or None if empty/unreadable.

    Zip-bomb defense: caps total file count and decompressed size (config).
    Path-traversal defense: rejects members with absolute paths or `..`.
    Sync — caller wraps in asyncio.to_thread.
    """
    from app.config import settings
    import zipfile
    import tarfile
    from io import BytesIO

    ext = ext.lower().lstrip(".")
    max_files = settings.archive_max_files
    max_bytes = settings.archive_max_total_mb * 1024 * 1024

    members: list[tuple[str, bytes]] = []
    total = 0

    def _safe_name(name: str) -> str | None:
        # Reject absolute paths and traversal — never extract outside the root.
        if name.startswith("/") or ".." in name.split("/"):
            return None
        return name

    try:
        if ext == "zip":
            with zipfile.ZipFile(BytesIO(raw)) as zf:
                for info in zf.infolist():
                    if info.is_dir():
                        continue
                    name = _safe_name(info.filename)
                    if not name:
                        continue
                    if len(members) >= max_files:
                        break
                    data = zf.read(info)
                    total += len(data)
                    if total > max_bytes:
                        break
                    members.append((name, data))
        elif ext in ("tar", "gz", "tgz"):
            mode = "r:gz" if ext in ("gz", "tgz") else "r:"
            with tarfile.open(fileobj=BytesIO(raw), mode=mode) as tf:
                for info in tf:
                    if not info.isfile():
                        continue
                    name = _safe_name(info.name)
                    if not name:
                        continue
                    if len(members) >= max_files:
                        break
                    f = tf.extractfile(info)
                    if f is None:
                        continue
                    data = f.read()
                    total += len(data)
                    if total > max_bytes:
                        break
                    members.append((name, data))
        else:
            return None
    except Exception:  # noqa: BLE001 — corrupt/unsupported archive
        return None

    if not members:
        return None

    parts: list[str] = []
    for name, data in members:
        m_ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
        if m_ext in OFFICE_EXTRACTORS or m_ext == "pdf":
            text = _extract_doc_content_sync(data, m_ext)
        elif m_ext in PLAIN_TEXT_EXTS:
            text = _decode_text(data)
        else:
            continue  # skip binaries inside archives
        if text and text.strip():
            parts.append(f"=== {name} ===\n{text}")
    return "\n\n".join(parts) if parts else None


def _extract_doc_content_sync(raw: bytes, ext: str) -> str | None:
    """Sync version of _extract_doc_content for use inside _extract_archive
    (which is already running in a thread)."""
    ext = ext.lower().lstrip(".")
    if ext in OFFICE_EXTRACTORS:
        return OFFICE_EXTRACTORS[ext](raw)
    if ext == "pdf":
        return extract_pdf_text(raw)
    if ext in PLAIN_TEXT_EXTS:
        return _decode_text(raw)
    return None


async def _extract_doc_content(
    raw: bytes, ext: str, *, prefer_docling: bool = True,
) -> str | None:
    """Unified document extraction: Docling first (Markdown + tables + OCR),
    falling back to the legacy per-format extractors.

    Docling handles pdf/docx/pptx/html. xlsx/csv stay on openpyxl (Docling is
    weak on pure data tables). Returns the extracted text, or a fallback error
    placeholder for Office, or None for PDF failures."""
    ext = ext.lower().lstrip(".")
    # Docling path: only for formats it handles well.
    if prefer_docling:
        from app.core.docling_converter import is_supported, convert_bytes_to_markdown_sync
        if is_supported(ext):
            md = await asyncio.to_thread(convert_bytes_to_markdown_sync, raw, ext)
            if md:
                return md
            # else: Docling failed or unavailable — fall through to legacy.
    # Legacy extractors.
    if ext in OFFICE_EXTRACTORS:
        return await asyncio.to_thread(OFFICE_EXTRACTORS[ext], raw) or "<p><em>(无法解析文档内容)</em></p>"
    if ext == "pdf":
        return await asyncio.to_thread(extract_pdf_text, raw)
    if ext in PLAIN_TEXT_EXTS:
        return _decode_text(raw)
    return None


async def process_upload(
    raw: bytes,
    ext: str,
    storage_key_prefix: str,
    name: str,
    content_type: str | None = None,
    fast_mode: bool = False,
) -> ProcessedUpload:
    """Decide how to store + extract an uploaded file's bytes.

    Single source of truth for "large vs. small file" handling, shared by
    every upload endpoint (conversation attachments, personal file storage,
    team knowledge base, project docs) so they no longer each reinvent (and
    subtly diverge on) this decision:

    - Office docs (docx/xlsx/pptx/csv/rtf): `content` holds extracted preview
      HTML, not the raw bytes, so the raw bytes always go to object storage
      regardless of size — otherwise the "download original" route would
      have nothing but HTML to serve back.
    - Anything else bigger than settings.file_offload_threshold_kb, or when
      the storage backend is minio: raw bytes go to object storage; text or
      PDF content is still extracted (best-effort) for prompt injection.
    - Everything else (small, non-office): inlined directly — text types
      decoded as-is, PDFs text-extracted, everything else base64.

    `storage_key_prefix` is the caller's namespace (e.g. "conversations/{id}"
    or "team-knowledge/{id}"); the object key becomes
    "{storage_key_prefix}/{uuid}/{name}" so the original filename/extension
    stays visible when browsing the bucket directly.
    """
    from app.config import settings
    from app.core import object_storage
    from app.core.file_validation import validate_upload

    ext = ext.lower()
    # P2-security: cross-check the declared extension against magic bytes so a
    # renamed executable can't slip through. No-op if python-magic is missing.
    validate_upload(raw, ext)
    ctype = content_type or "application/octet-stream"
    threshold_bytes = settings.file_offload_threshold_kb * 1024
    storage_key: str | None = None
    content: str | None = None
    preview_pdf_key: str | None = None

    # P2-file: strip EXIF from images (privacy: GPS/camera metadata).
    if settings.strip_exif_enabled and ext in ("jpg", "jpeg", "png", "webp"):
        raw = await asyncio.to_thread(_strip_exif, raw, ext)

    # P2-file: archives (zip/tar/gz) — extract and concatenate member text.
    if ext in ("zip", "tar", "gz", "tgz"):
        storage_key = f"{storage_key_prefix}/{uuid.uuid4().hex}/{name}"
        await asyncio.to_thread(object_storage.put, storage_key, raw, ctype)
        content = await asyncio.to_thread(_extract_archive, raw, ext)
        return ProcessedUpload(content=content, storage_key=storage_key, size_bytes=len(raw))

    if ext in OFFICE_EXTRACTORS:
        storage_key = f"{storage_key_prefix}/{uuid.uuid4().hex}/{name}"
        try:
            await asyncio.to_thread(object_storage.put, storage_key, raw, ctype)
        except Exception as exc:
            raise HTTPException(
                status_code=503, detail="文件预览服务不可用，请检查对象存储配置"
            ) from exc
        # `content` always comes from the per-format extractor (csv/xlsx keep
        # their structured HTML tables, rtf its text, docx the mammoth HTML,
        # legacy formats the soffice HTML) — NEVER from PDF text extraction,
        # which would destroy table structure and waste a soffice conversion.
        # LibreOffice PDF conversion is preview-only: it fills
        # preview_pdf_key and nothing else. csv/rtf skip it entirely (they
        # have no PDF preview in the UI).
        content = await _extract_doc_content(raw, ext, prefer_docling=not fast_mode)
        if ext in _PDF_PREVIEW_EXTS:
            from app.core.office_legacy import convert_to_pdf
            pdf_bytes = await asyncio.to_thread(convert_to_pdf, raw, ext)
            if pdf_bytes:
                preview_pdf_key = f"{storage_key_prefix}/previews/{uuid.uuid4().hex}.pdf"
                try:
                    await asyncio.to_thread(
                        object_storage.put, preview_pdf_key, pdf_bytes, "application/pdf"
                    )
                except Exception:
                    logger.warning("failed to store preview pdf for .%s upload", ext, exc_info=True)
                    preview_pdf_key = None
    elif len(raw) > threshold_bytes or settings.storage_backend == "minio":
        storage_key = f"{storage_key_prefix}/{uuid.uuid4().hex}/{name}"
        await asyncio.to_thread(object_storage.put, storage_key, raw, ctype)
        if ext == "pdf":
            content = await _extract_doc_content(raw, ext, prefer_docling=not fast_mode)
        elif ext in PLAIN_TEXT_EXTS:
            content = _decode_text(raw)
    else:
        if ext in PLAIN_TEXT_EXTS:
            content = _decode_text(raw)
        elif ext == "pdf":
            content = await _extract_doc_content(raw, ext, prefer_docling=not fast_mode)
        else:
            content = base64.b64encode(raw).decode("ascii")

    return ProcessedUpload(
        content=content, storage_key=storage_key, size_bytes=len(raw),
        preview_pdf_key=preview_pdf_key,
    )


async def hydrate_stored_content(
    kind: str | None,
    storage_key: str | None,
    inline_content: str | None,
) -> str | None:
    """Resolve knowledge/doc content from inline text or object storage.

    Consolidates the "if content is None and storage_key: fetch from storage
    and extract" pattern that was duplicated across teams.py and
    conversation_service.py. Returns ``None`` if no content is available.
    """
    if inline_content is not None:
        return inline_content
    if not storage_key or not kind:
        return None

    from app.core import object_storage

    data = await asyncio.to_thread(object_storage.get, storage_key)
    ext = kind.lower()
    if ext in OFFICE_EXTRACTORS:
        # to_thread: legacy extractors may spawn a soffice subprocess (45s
        # timeout cap) — running it on the event loop would freeze the whole
        # API for every request touching an un-extracted legacy file.
        return await asyncio.to_thread(OFFICE_EXTRACTORS[ext], data) or None
    if is_text_extractable(ext):
        return data.decode("utf-8", "ignore")
    return None


async def read_upload_capped(file: UploadFile, max_bytes: int) -> bytes:
    """Read an UploadFile fully, but abort with HTTP 413 once it exceeds max_bytes.

    Reads in chunks so an oversized upload can't balloon memory before the
    limit is hit (``await file.read()`` would buffer the whole body first).
    """
    chunks: list[bytes] = []
    total = 0
    while True:
        chunk = await file.read(_UPLOAD_CHUNK)
        if not chunk:
            break
        total += len(chunk)
        if total > max_bytes:
            raise HTTPException(
                status_code=413,
                detail=f"文件过大，上限 {max_bytes // (1024 * 1024)}MB",
            )
        chunks.append(chunk)
    return b"".join(chunks)


def safe_relative_path(name: str, fallback: str = "untitled.txt") -> str:
    """Normalize a user/agent-supplied path to a contained relative path.
    Anchors the path at root before normalizing so ``../`` segments can never
    climb above it (``a/../../b`` → ``b``, ``../../etc/passwd`` → ``etc/passwd``),
    then strips the leading separator. Never raises — preserves valid nested
    paths like ``src/main.py`` for folder support.
    """
    candidate = (name or "").replace("\\", "/").strip()
    normalized = os.path.normpath("/" + candidate).lstrip("/")
    return normalized.replace(os.sep, "/") or fallback


def safe_download_headers(name: str, mime: str) -> dict[str, str]:
    """Response headers for file downloads/previews.

    SVG is served as an ATTACHMENT with a sandbox CSP: an uploaded SVG can
    embed <script>/onload payloads, and serving it inline from the app origin
    would be stored XSS (the payload runs with our origin's localStorage /
    cookies). <img> embedding is safe, direct navigation is not.
    """
    from urllib.parse import quote
    ascii_name = (name or "file").encode("ascii", "ignore").decode() or "file"
    disposition = "attachment" if mime == "image/svg+xml" else "inline"
    headers = {
        "Content-Disposition": (
            f'{disposition}; filename="{ascii_name}"; filename*=UTF-8\'\'{quote(name or "file")}'
        ),
    }
    if mime == "image/svg+xml":
        headers["Content-Security-Policy"] = "sandbox"
    return headers


def confine_to_dir(base_dir: str, relative: str) -> str:
    """Join base_dir + a (pre-normalized) relative path and assert containment.

    Defense in depth after ``safe_relative_path``: resolves symlinks and rejects
    any result that escapes base_dir. Raises ValueError on escape — callers in
    the service layer translate it to a skip/400 as appropriate (an HTTP
    exception here would couple the files module to the API layer).
    """
    base_real = os.path.realpath(base_dir)
    target = os.path.realpath(os.path.join(base_real, relative))
    if target != base_real and not target.startswith(base_real + os.sep):
        raise ValueError(f"path escapes base dir: {relative}")
    return target
