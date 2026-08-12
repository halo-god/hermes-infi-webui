"""Office document (docx/xlsx/pptx) -> HTML preview extraction.

Builds minimal fixtures in-memory with the same libraries used to extract,
round-trips them through each extractor, and asserts the XSS guard (embedded
literal "<script>" text must come out escaped, never as an executable tag).
"""
import io
import uuid

import pytest
from httpx import AsyncClient

from app.core.files import (
    _XLSX_MAX_COLS,
    _XLSX_MAX_ROWS,
    extract_docx_html,
    extract_pptx_html,
    extract_xlsx_html,
)


def _docx_bytes() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("Title", level=1)
    doc.add_paragraph("A plain paragraph.")
    p = doc.add_paragraph()
    p.add_run("bold text").bold = True
    doc.add_paragraph("<script>alert(1)</script>")
    doc.add_paragraph("Bullet one", style="List Bullet")
    doc.add_paragraph("Bullet two", style="List Bullet")
    doc.add_paragraph("Numbered one", style="List Number")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Header A"
    table.rows[0].cells[1].text = "Header B"
    table.rows[1].cells[0].text = "cell 1"
    table.rows[1].cells[1].text = "cell 2"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_extract_docx_html_structure_and_escaping():
    html = extract_docx_html(_docx_bytes())
    assert html is not None
    assert "<h1>Title</h1>" in html
    assert "<p>A plain paragraph.</p>" in html
    assert "<strong>bold text</strong>" in html
    assert "<ul>" in html and "<li>Bullet one</li>" in html and "<li>Bullet two</li>" in html
    assert "<ol>" in html and "<li>Numbered one</li>" in html
    # Mammoth emits semantic HTML: cells are <td> wrapping a <p> (no th
    # distinction), text still present and escaped.
    assert "<table>" in html
    assert "Header A" in html and "cell 1" in html
    # XSS guard: literal script text must be escaped, never an executable tag.
    assert "<script>" not in html
    assert "&lt;script&gt;" in html


def test_extract_docx_html_handles_garbage_input():
    assert extract_docx_html(b"not a real docx") is None


def _xlsx_bytes(rows: int = 3, cols: int = 3) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet<script>"
    for r in range(1, rows + 1):
        for c in range(1, cols + 1):
            ws.cell(row=r, column=c, value=f"r{r}c{c}")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_extract_xlsx_html_basic():
    html = extract_xlsx_html(_xlsx_bytes())
    assert html is not None
    assert "&lt;script&gt;" in html  # sheet name escaped
    assert "<script>" not in html
    assert "<table>" in html and "<td>r1c1</td>" in html


def test_extract_xlsx_html_truncates_oversized_sheets():
    html = extract_xlsx_html(_xlsx_bytes(rows=_XLSX_MAX_ROWS + 10, cols=_XLSX_MAX_COLS + 5))
    assert html is not None
    assert f"仅显示前 {_XLSX_MAX_ROWS} 行" in html
    assert f"仅显示前 {_XLSX_MAX_COLS} 列" in html


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Slide One Title"
    s1.placeholders[1].text = "<script>alert(1)</script>"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Slide Two Title"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_html_per_slide():
    html = extract_pptx_html(_pptx_bytes())
    assert html is not None
    assert "Slide 1" in html and "Slide 2" in html
    assert "Slide One Title" in html and "Slide Two Title" in html
    assert "<script>" not in html
    assert "&lt;script&gt;" in html


@pytest.mark.asyncio
async def test_upload_docx_preview_and_raw_download_roundtrip(client: AsyncClient, auth_headers, test_user, db):
    """Uploading a real .docx must: accept immediately with
    processing_status="processing" (conversion is async), then end with
    extractable text in `content` (PDF-extracted when soffice is present,
    HTML otherwise), keep the original bytes byte-identical via the
    raw-download route (the dual-storage split), and tag `kind` correctly.
    With soffice installed the LibreOffice PDF preview key is set and /pdf
    serves a real PDF."""
    pytest.importorskip("moto")
    from moto import mock_aws

    from app.config import settings
    from app.core import object_storage
    from app.db.models.conversation import Conversation

    convo = Conversation(owner_id=test_user.id, title="office-upload-test")
    db.add(convo)
    await db.commit()
    await db.refresh(convo)

    raw = _docx_bytes()

    prev_endpoint = settings.minio_endpoint
    prev_bucket = settings.minio_bucket
    settings.minio_endpoint = ""
    settings.minio_bucket = "hermes-test-bucket"

    # moto's mock_aws must be used as a context manager (not a decorator)
    # around async code — as a decorator on an async function it tears the
    # mock down before the coroutine actually runs, since calling an async
    # function just returns a coroutine object immediately.
    try:
        with mock_aws():
            object_storage.reset_client()
            r = await client.post(
                f"/api/v1/conversations/{convo.id}/upload",
                files={"file": ("report.docx", raw, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
                headers=auth_headers,
            )
            assert r.status_code == 201, r.text
            data = r.json()
            assert data["kind"] == "docx"
            # Async conversion: the upload returns immediately in "processing".
            assert data["processing_status"] == "processing"
            file_id = data["id"]

            # The upload endpoint accepts instantly; the conversion runs in a
            # background task. That task opens its own session and cannot see
            # rows still inside this test's transaction (the db fixture wraps
            # everything in a rolled-back transaction), so simulate its effect
            # here: run process_upload and flip the row to ready.
            from app.core.files import process_upload
            from app.db.models.workspace import WorkspaceFile
            wf = await db.get(WorkspaceFile, uuid.UUID(file_id))
            assert wf is not None
            processed = await process_upload(
                raw, "docx", f"conversations/{convo.id}", "report.docx",
                content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                fast_mode=True,
            )
            wf.content = processed.content
            wf.storage_key = processed.storage_key
            wf.preview_pdf_key = processed.preview_pdf_key
            wf.processing_status = "ready"
            await db.commit()
            # updated_at is a server-side onupdate column — after commit it is
            # stale-on-this-instance; refresh so the shared fixture session
            # (identity-mapped to the same object) doesn't lazy-load it inside
            # the sync pydantic validation of the detail endpoint.
            await db.refresh(wf)

            detail = await client.get(f"/api/v1/conversations/{convo.id}/files/{file_id}", headers=auth_headers)
            assert detail.status_code == 200
            assert detail.json()["processing_status"] == "ready"
            content = detail.json()["content"]
            # "Title" survives both pipelines: mammoth HTML fallback
            # (<h1>Title</h1>) and PDF text extraction.
            assert "Title" in content
            # With soffice installed, the unified PDF preview is produced and
            # served by the /pdf endpoint.
            pdf_key = detail.json().get("preview_pdf_key")
            if pdf_key:
                pdf_resp = await client.get(f"/api/v1/conversations/{convo.id}/files/{file_id}/pdf", headers=auth_headers)
                assert pdf_resp.status_code == 200
                assert pdf_resp.content[:4] == b"%PDF"

            raw_resp = await client.get(f"/api/v1/conversations/{convo.id}/files/{file_id}/raw", headers=auth_headers)
            assert raw_resp.status_code == 200
            assert raw_resp.content == raw  # byte-identical original, not the preview
    finally:
        settings.minio_endpoint = prev_endpoint
        settings.minio_bucket = prev_bucket
        object_storage.reset_client()


# ── Legacy formats (soffice optional) ──


def test_legacy_doc_hint_when_soffice_absent(monkeypatch):
    """Without LibreOffice, legacy formats yield a clear hint (not silent)."""
    from app.core.office_legacy import extract_legacy_doc_html

    monkeypatch.setattr("app.core.office_legacy.soffice_available", lambda: False)
    hint = extract_legacy_doc_html(b"not a real doc", "doc")
    assert hint is not None
    assert "暂不支持" in hint and "docx" in hint


def test_legacy_csv_table_wrapping():
    from app.core.office_legacy import _csv_to_table

    tbl = _csv_to_table('a,b\n1,2\n')
    assert "<table>" in tbl
    assert "<th>a</th>" in tbl and "<td>1</td>" in tbl
    assert "<script>" not in tbl  # cells escaped


def test_legacy_formats_registered_in_extractors():
    from app.core.files import OFFICE_EXTRACTORS

    for ext in ("doc", "xls", "ppt"):
        assert ext in OFFICE_EXTRACTORS, f"{ext} must be registered"
